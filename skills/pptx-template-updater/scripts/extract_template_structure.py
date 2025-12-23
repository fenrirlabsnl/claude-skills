#!/usr/bin/env python3
"""
Extract semantic structure from a PowerPoint template.

This script analyzes a PPTX file and outputs detailed information about each shape
including text content, position, type, and character count. This helps derive
semantic meaning from the template structure.

Usage:
    python extract_template_structure.py <template.pptx> [--output output.json]
"""

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError, PythonPptxError

# Import security utilities
from security_utils import (
    validate_input_file,
    validate_output_file,
    SecurityError,
    PathTraversalError,
    FileSizeError
)


def get_shape_type_name(shape_type):
    """Convert shape type enum to readable name."""
    type_names = {
        MSO_SHAPE_TYPE.AUTO_SHAPE: "AutoShape",
        MSO_SHAPE_TYPE.CALLOUT: "Callout",
        MSO_SHAPE_TYPE.CHART: "Chart",
        MSO_SHAPE_TYPE.COMMENT: "Comment",
        MSO_SHAPE_TYPE.FREEFORM: "Freeform",
        MSO_SHAPE_TYPE.GROUP: "Group",
        MSO_SHAPE_TYPE.LINE: "Line",
        MSO_SHAPE_TYPE.MEDIA: "Media",
        MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT: "OLEControl",
        MSO_SHAPE_TYPE.PICTURE: "Picture",
        MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
        MSO_SHAPE_TYPE.TABLE: "Table",
        MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
    }
    return type_names.get(shape_type, "Unknown")


def extract_text_from_shape(shape):
    """Extract all text content from a shape, including nested text frames."""
    if not hasattr(shape, "text_frame"):
        return ""

    if not shape.has_text_frame:
        return ""

    return shape.text_frame.text


def count_paragraphs_and_bullets(shape):
    """Count paragraphs and bullet points in a shape."""
    if not hasattr(shape, "text_frame") or not shape.has_text_frame:
        return {"paragraphs": 0, "bullets": 0}

    paragraphs = len(shape.text_frame.paragraphs)
    bullets = sum(1 for p in shape.text_frame.paragraphs if p.level > 0 or p.text.strip())

    return {"paragraphs": paragraphs, "bullets": bullets}


def extract_table_data(table):
    """Extract data from a table shape."""
    table_data = {
        "rows": len(table.rows),
        "columns": len(table.columns),
        "cells": []
    }

    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            cell_text = cell.text_frame.text if cell.text_frame else ""
            if cell_text:  # Only include cells with content
                table_data["cells"].append({
                    "row": row_idx,
                    "column": col_idx,
                    "text": cell_text,
                    "character_count": len(cell_text)
                })

    return table_data


def analyze_shape(shape, shape_index):
    """Analyze a single shape and extract relevant metadata."""
    shape_data = {
        "index": shape_index,
        "name": shape.name,
        "shape_type": get_shape_type_name(shape.shape_type),
        "position": {
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height,
        },
    }

    # Handle tables separately
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        table_data = extract_table_data(shape.table)
        shape_data["is_table"] = True
        shape_data["table"] = table_data
        shape_data["text_content"] = ""  # Tables don't have direct text content
        shape_data["character_count"] = sum(cell["character_count"] for cell in table_data["cells"])
        shape_data["paragraphs"] = 0
        shape_data["bullets"] = 0
    else:
        # Regular shape with text frame
        text_content = extract_text_from_shape(shape)
        para_info = count_paragraphs_and_bullets(shape)

        shape_data["text_content"] = text_content
        shape_data["character_count"] = len(text_content)
        shape_data["paragraphs"] = para_info["paragraphs"]
        shape_data["bullets"] = para_info["bullets"]
        shape_data["is_table"] = False

    # Add placeholder info if applicable
    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
        shape_data["is_placeholder"] = True
        shape_data["placeholder_type"] = str(shape.placeholder_format.type)

    return shape_data


def extract_template_structure(pptx_path):
    """Extract complete structure from a PowerPoint template."""
    prs = Presentation(pptx_path)

    structure = {
        "file_name": Path(pptx_path).name,
        "total_slides": len(prs.slides),
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": []
    }

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": slide_index,
            "shapes": []
        }

        for shape_index, shape in enumerate(slide.shapes, start=1):
            try:
                shape_data = analyze_shape(shape, shape_index)
                # Include shapes with text OR tables with content
                if shape_data["text_content"] or (shape_data.get("is_table") and shape_data.get("table", {}).get("cells")):
                    slide_data["shapes"].append(shape_data)
            except (AttributeError, KeyError, IndexError) as e:
                print(f"Warning: Could not analyze shape {shape_index} on slide {slide_index}: {e}",
                      file=sys.stderr)
            except PythonPptxError as e:
                print(f"Warning: PowerPoint error analyzing shape {shape_index} on slide {slide_index}: {e}",
                      file=sys.stderr)

        structure["slides"].append(slide_data)

    return structure


def main():
    parser = argparse.ArgumentParser(
        description="Extract semantic structure from a PowerPoint template"
    )
    parser.add_argument("template", help="Path to the PPTX template file")
    parser.add_argument(
        "--output", "-o",
        help="Output JSON file (default: print to stdout)",
        default=None
    )

    args = parser.parse_args()

    # Validate inputs with security checks
    try:
        # Validate template file (PPTX)
        template_path = validate_input_file(
            args.template,
            allowed_extensions=['.pptx']
        )

        # Validate output file path if specified
        output_path = None
        if args.output:
            output_path = validate_output_file(
                args.output,
                allowed_extensions=['.json']
            )

    except PathTraversalError as e:
        print(f"Security Error: {e}", file=sys.stderr)
        sys.exit(1)
    except FileSizeError as e:
        print(f"File Size Error: {e}", file=sys.stderr)
        sys.exit(1)
    except FileNotFoundError as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
    except ValueError as e:
        print(f"Invalid file type: {e}", file=sys.stderr)
        sys.exit(1)
    except SecurityError as e:
        print(f"Security Error: {e}", file=sys.stderr)
        sys.exit(1)

    # Extract template structure
    try:
        structure = extract_template_structure(template_path)
        output_json = json.dumps(structure, indent=2)

        if output_path:
            output_path.write_text(output_json, encoding='utf-8')
            print(f"Structure extracted to: {output_path}", file=sys.stderr)
        else:
            print(output_json)

    except PackageNotFoundError as e:
        print(f"Error: Invalid or corrupted PowerPoint file: {e}", file=sys.stderr)
        sys.exit(1)
    except PermissionError as e:
        print(f"Error: Permission denied: {e}", file=sys.stderr)
        sys.exit(1)
    except OSError as e:
        print(f"Error: File system error: {e}", file=sys.stderr)
        sys.exit(1)
    except PythonPptxError as e:
        print(f"Error processing PowerPoint file: {e}", file=sys.stderr)
        sys.exit(1)
    except UnicodeEncodeError as e:
        print(f"Error: Character encoding error in output: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
