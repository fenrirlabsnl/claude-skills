#!/usr/bin/env python3
"""
Update a PowerPoint template with new content while preserving structure.

This script updates text in PowerPoint shapes based on update instructions,
preserving formatting, bullet structures, and respecting length constraints.

Usage:
    python update_template.py <template.pptx> <updates.json> <output.pptx>

Updates JSON format:
{
  "updates": [
    {
      "slide": 1,
      "shape": 2,
      "text": "New content here",
      "preserve_bullets": true
    }
  ]
}
"""

import sys
import json
import argparse
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def update_table_cell(table, row, col, new_text):
    """
    Update a specific cell in a table while preserving formatting.

    Args:
        table: The table object
        row: Row index (0-based)
        col: Column index (0-based)
        new_text: New text for the cell

    Returns:
        dict: Update result
    """
    try:
        cell = table.rows[row].cells[col]
        text_frame = cell.text_frame

        # Capture original formatting
        original_font_size = None
        original_font_name = None
        original_font_bold = None
        original_font_italic = None
        original_font_underline = None
        original_font_color = None

        if text_frame.paragraphs:
            first_para = text_frame.paragraphs[0]

            # Try to get formatting from first run
            if first_para.runs:
                first_run = first_para.runs[0]
                original_font_size = first_run.font.size
                original_font_name = first_run.font.name
                original_font_bold = first_run.font.bold
                original_font_italic = first_run.font.italic
                original_font_underline = first_run.font.underline
                # Store color object (handles RGB and theme colors)
                if first_run.font.color.type is not None:
                    original_font_color = first_run.font.color

            # Fallback to paragraph's font properties if no runs exist
            if original_font_size is None:
                original_font_size = first_para.font.size
            if original_font_name is None:
                original_font_name = first_para.font.name
            if original_font_bold is None:
                original_font_bold = first_para.font.bold
            if original_font_italic is None:
                original_font_italic = first_para.font.italic
            if original_font_underline is None:
                original_font_underline = first_para.font.underline

        # Clear and update cell
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = new_text

        # Apply original formatting
        if original_font_size:
            run.font.size = original_font_size
        if original_font_name:
            run.font.name = original_font_name
        if original_font_bold is not None:
            run.font.bold = original_font_bold
        if original_font_italic is not None:
            run.font.italic = original_font_italic
        if original_font_underline is not None:
            run.font.underline = original_font_underline
        if original_font_color is not None:
            run.font.color.rgb = original_font_color

        return {"success": True, "cell": f"({row},{col})"}

    except Exception as e:
        return {"success": False, "error": str(e), "cell": f"({row},{col})"}


def update_shape_text(shape, new_text, preserve_bullets=True, warn_on_overflow=True):
    """
    Update shape text while preserving formatting.

    Args:
        shape: The shape to update
        new_text: New text content (can include newlines for bullets)
        preserve_bullets: Whether to maintain bullet point structure
        warn_on_overflow: Warn if new text is significantly longer than original

    Returns:
        dict: Update result with warnings
    """
    if not hasattr(shape, "text_frame") or not shape.has_text_frame:
        return {"success": False, "error": "Shape has no text frame"}

    original_length = len(shape.text_frame.text)
    result = {"success": True, "warnings": []}

    # Check length
    if warn_on_overflow and len(new_text) > original_length * 1.5:
        result["warnings"].append(
            f"New text ({len(new_text)} chars) is significantly longer than "
            f"original ({original_length} chars). May overflow shape."
        )

    text_frame = shape.text_frame

    # Capture original formatting from first run or paragraph defaults
    original_font_size = None
    original_font_name = None
    original_font_bold = None
    original_font_italic = None
    original_font_underline = None
    original_font_color = None

    if text_frame.paragraphs:
        first_para = text_frame.paragraphs[0]

        # Try to get formatting from first run
        if first_para.runs:
            first_run = first_para.runs[0]
            original_font_size = first_run.font.size
            original_font_name = first_run.font.name
            original_font_bold = first_run.font.bold
            original_font_italic = first_run.font.italic
            original_font_underline = first_run.font.underline
            # Store color object (handles RGB and theme colors)
            if first_run.font.color.type is not None:
                original_font_color = first_run.font.color

        # Fallback to paragraph's font properties if no runs exist
        if original_font_size is None:
            original_font_size = first_para.font.size
        if original_font_name is None:
            original_font_name = first_para.font.name
        if original_font_bold is None:
            original_font_bold = first_para.font.bold
        if original_font_italic is None:
            original_font_italic = first_para.font.italic
        if original_font_underline is None:
            original_font_underline = first_para.font.underline

    if preserve_bullets and '\n' in new_text:
        # Preserve bullet structure
        lines = new_text.split('\n')

        # Clear existing paragraphs but keep the first one
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]
            p._element.getparent().remove(p._element)

        # Update first paragraph with formatting preservation
        if lines:
            first_para = text_frame.paragraphs[0]
            first_para.clear()
            run = first_para.add_run()
            run.text = lines[0]
            # Apply original formatting
            if original_font_size:
                run.font.size = original_font_size
            if original_font_name:
                run.font.name = original_font_name
            if original_font_bold is not None:
                run.font.bold = original_font_bold
            if original_font_italic is not None:
                run.font.italic = original_font_italic
            if original_font_underline is not None:
                run.font.underline = original_font_underline
            if original_font_color is not None:
                # Copy color (handles RGB and theme colors)
                try:
                    if hasattr(original_font_color, 'rgb') and original_font_color.rgb:
                        run.font.color.rgb = original_font_color.rgb
                    elif hasattr(original_font_color, 'theme_color'):
                        run.font.color.theme_color = original_font_color.theme_color
                except:
                    pass  # Skip if color can't be applied

        # Add remaining paragraphs with formatting
        for line in lines[1:]:
            p = text_frame.add_paragraph()
            run = p.add_run()
            run.text = line
            # Copy level from first paragraph
            p.level = text_frame.paragraphs[0].level
            # Apply original formatting
            if original_font_size:
                run.font.size = original_font_size
            if original_font_name:
                run.font.name = original_font_name
            if original_font_bold is not None:
                run.font.bold = original_font_bold
            if original_font_italic is not None:
                run.font.italic = original_font_italic
            if original_font_underline is not None:
                run.font.underline = original_font_underline
            if original_font_color is not None:
                # Copy color (handles RGB and theme colors)
                try:
                    if hasattr(original_font_color, 'rgb') and original_font_color.rgb:
                        run.font.color.rgb = original_font_color.rgb
                    elif hasattr(original_font_color, 'theme_color'):
                        run.font.color.theme_color = original_font_color.theme_color
                except:
                    pass  # Skip if color can't be applied

    else:
        # Simple text replacement with formatting preservation
        first_para = text_frame.paragraphs[0]
        first_para.clear()
        run = first_para.add_run()
        run.text = new_text
        # Apply original formatting
        if original_font_size:
            run.font.size = original_font_size
        if original_font_name:
            run.font.name = original_font_name
        if original_font_bold is not None:
            run.font.bold = original_font_bold
        if original_font_italic is not None:
            run.font.italic = original_font_italic
        if original_font_underline is not None:
            run.font.underline = original_font_underline
        if original_font_color is not None:
            run.font.color.rgb = original_font_color

    result["new_length"] = len(new_text)
    result["original_length"] = original_length
    return result


def apply_updates(pptx_path, updates_data, output_path):
    """
    Apply updates to a PowerPoint template.

    Args:
        pptx_path: Path to input PPTX template
        updates_data: Dictionary with update instructions
        output_path: Path for output PPTX file

    Returns:
        dict: Results of all updates
    """
    prs = Presentation(pptx_path)
    results = {"updates_applied": 0, "errors": [], "warnings": []}

    for update in updates_data.get("updates", []):
        slide_num = update.get("slide")
        shape_idx = update.get("shape")
        new_text = update.get("text", "")
        preserve_bullets = update.get("preserve_bullets", True)
        table_cells = update.get("table_cells", [])  # For table updates

        # Validate slide number
        if slide_num < 1 or slide_num > len(prs.slides):
            results["errors"].append(
                f"Invalid slide number: {slide_num}. Template has {len(prs.slides)} slides."
            )
            continue

        slide = prs.slides[slide_num - 1]

        # Validate shape index
        if shape_idx < 1 or shape_idx > len(slide.shapes):
            results["errors"].append(
                f"Invalid shape index: {shape_idx} on slide {slide_num}. "
                f"Slide has {len(slide.shapes)} shapes."
            )
            continue

        shape = slide.shapes[shape_idx - 1]

        # Check if this is a table update
        if table_cells and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Update table cells
            table = shape.table
            for cell_update in table_cells:
                row = cell_update.get("row")
                col = cell_update.get("column")
                cell_text = cell_update.get("text", "")

                cell_result = update_table_cell(table, row, col, cell_text)

                if cell_result["success"]:
                    results["updates_applied"] += 1
                else:
                    results["errors"].append(
                        f"Slide {slide_num}, Shape {shape_idx}, Cell {cell_result['cell']}: {cell_result.get('error')}"
                    )

        else:
            # Apply regular text update
            update_result = update_shape_text(
                shape,
                new_text,
                preserve_bullets=preserve_bullets
            )

            if update_result["success"]:
                results["updates_applied"] += 1
                if update_result.get("warnings"):
                    for warning in update_result["warnings"]:
                        results["warnings"].append(
                            f"Slide {slide_num}, Shape {shape_idx}: {warning}"
                        )
            else:
                results["errors"].append(
                    f"Slide {slide_num}, Shape {shape_idx}: {update_result.get('error')}"
                )

    # Save the presentation
    prs.save(output_path)
    results["output_file"] = str(output_path)

    return results


def main():
    parser = argparse.ArgumentParser(
        description="Update PowerPoint template with new content"
    )
    parser.add_argument("template", help="Path to the input PPTX template")
    parser.add_argument("updates", help="Path to JSON file with updates")
    parser.add_argument("output", help="Path for the output PPTX file")

    args = parser.parse_args()

    # Validate inputs
    if not Path(args.template).exists():
        print(f"Error: Template file not found: {args.template}", file=sys.stderr)
        sys.exit(1)

    if not Path(args.updates).exists():
        print(f"Error: Updates file not found: {args.updates}", file=sys.stderr)
        sys.exit(1)

    # Load updates
    try:
        with open(args.updates, 'r') as f:
            updates_data = json.load(f)
    except json.JSONDecodeError as e:
        print(f"Error: Invalid JSON in updates file: {e}", file=sys.stderr)
        sys.exit(1)

    # Apply updates
    try:
        results = apply_updates(args.template, updates_data, args.output)

        # Print results
        print(f"\n✅ Updates applied: {results['updates_applied']}")
        print(f"📄 Output saved to: {results['output_file']}")

        if results["warnings"]:
            print(f"\n⚠️  Warnings ({len(results['warnings'])}):")
            for warning in results["warnings"]:
                print(f"  - {warning}")

        if results["errors"]:
            print(f"\n❌ Errors ({len(results['errors'])}):")
            for error in results["errors"]:
                print(f"  - {error}")

    except Exception as e:
        print(f"Error updating template: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
